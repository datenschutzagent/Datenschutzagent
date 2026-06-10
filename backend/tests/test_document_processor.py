"""Unit tests for document text extraction logic (no DB, no external services)."""
import io

import pytest

import app.services.document_processor as dp
from app.services.document_processor import (
    extract_text_from_docx,
    extract_text_from_xlsx,
    extract_text,
    ExtractionResult,
    EXTRACTION_METHOD_TEXT,
)


def _make_docx_bytes(paragraphs: list[str], table_rows: list[list[str]] | None = None) -> bytes:
    """Build a minimal in-memory DOCX file."""
    import docx as python_docx

    doc = python_docx.Document()
    for para in paragraphs:
        doc.add_paragraph(para)
    if table_rows:
        table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r_idx, row in enumerate(table_rows):
            for c_idx, cell_text in enumerate(row):
                table.rows[r_idx].cells[c_idx].text = cell_text
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_docx_with_header_footer(body: str, header: str, footer: str) -> bytes:
    """Build a DOCX whose header/footer carry text the body does not."""
    import docx as python_docx

    doc = python_docx.Document()
    doc.add_paragraph(body)
    section = doc.sections[0]
    section.header.paragraphs[0].text = header
    section.footer.paragraphs[0].text = footer
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes(sheets: dict[str, list[list]]) -> bytes:
    """Build a minimal in-memory XLSX file."""
    import openpyxl

    wb = openpyxl.Workbook()
    first = True
    for sheet_name, rows in sheets.items():
        if first:
            ws = wb.active
            ws.title = sheet_name
            first = False
        else:
            ws = wb.create_sheet(sheet_name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestExtractTextFromDocx:
    def test_paragraphs_extracted(self):
        content = _make_docx_bytes(["Erster Absatz", "Zweiter Absatz"])
        text = extract_text_from_docx(content)
        assert "Erster Absatz" in text
        assert "Zweiter Absatz" in text

    def test_table_cells_extracted(self):
        content = _make_docx_bytes([], table_rows=[["Name", "Wert"], ["Alice", "42"]])
        text = extract_text_from_docx(content)
        assert "Name" in text
        assert "Alice" in text

    def test_empty_document(self):
        content = _make_docx_bytes([])
        text = extract_text_from_docx(content)
        assert isinstance(text, str)

    def test_returns_string(self):
        content = _make_docx_bytes(["Hello"])
        result = extract_text_from_docx(content)
        assert isinstance(result, str)


class TestDocxHeaderFooter:
    def test_header_and_footer_text_recovered(self):
        # Header/footer often carry the controller, version/date or file reference, which plain
        # paragraph iteration drops. They must survive extraction.
        content = _make_docx_with_header_footer(
            body="Verarbeitungstätigkeit Lohnabrechnung",
            header="Verantwortlicher: Muster GmbH",
            footer="Stand: 01.01.2026 — Az. DS-2026-042",
        )
        text = extract_text_from_docx(content)
        assert "Verarbeitungstätigkeit Lohnabrechnung" in text
        assert "Verantwortlicher: Muster GmbH" in text
        assert "Az. DS-2026-042" in text


class TestExtractTextFromXlsx:
    def test_single_sheet_data(self):
        content = _make_xlsx_bytes({"Sheet1": [["Datum", "Betrag"], ["2024-01-01", 100]]})
        text = extract_text_from_xlsx(content)
        assert "Datum" in text
        assert "2024-01-01" in text

    def test_multiple_sheets(self):
        content = _make_xlsx_bytes({
            "Blatt1": [["A", "B"]],
            "Blatt2": [["C", "D"]],
        })
        text = extract_text_from_xlsx(content)
        assert "Blatt1" in text
        assert "Blatt2" in text
        assert "A" in text
        assert "C" in text

    def test_empty_sheet(self):
        content = _make_xlsx_bytes({"Empty": []})
        text = extract_text_from_xlsx(content)
        assert isinstance(text, str)


class TestExtractTextDispatcher:
    def test_docx_dispatch(self):
        content = _make_docx_bytes(["Testinhalt"])
        result = extract_text("file.docx", content)
        assert isinstance(result, ExtractionResult)
        assert result.extraction_method == EXTRACTION_METHOD_TEXT
        assert "Testinhalt" in result.text

    def test_xlsx_dispatch(self):
        content = _make_xlsx_bytes({"S1": [["Zelle"]]})
        result = extract_text("report.xlsx", content)
        assert isinstance(result, ExtractionResult)
        assert result.extraction_method == EXTRACTION_METHOD_TEXT
        assert "Zelle" in result.text

    def test_plain_text_fallback(self):
        content = "Hallo Welt".encode("utf-8")
        result = extract_text("readme.txt", content)
        assert isinstance(result, ExtractionResult)
        assert "Hallo Welt" in result.text

    def test_binary_fallback_returns_empty(self):
        content = b"\xff\xfe\x00\x01"  # not valid UTF-8
        result = extract_text("unknown.bin", content)
        assert isinstance(result, ExtractionResult)
        assert result.text == ""

    def test_cp1252_text_is_decoded_not_lost(self):
        # Legacy exports are often CP1252/Latin-1; decoding as UTF-8 fails. Such text must not be
        # silently dropped (no NUL bytes → it is real text, not binary).
        content = "Müller Straße – Zweck: Personalverwaltung".encode("cp1252")
        result = extract_text("export.txt", content)
        assert result.text != ""
        assert "ller" in result.text and "Personalverwaltung" in result.text

    def test_docx_uppercase_extension(self):
        content = _make_docx_bytes(["Groß"])
        result = extract_text("DOC.DOCX", content)
        assert "Groß" in result.text


class TestStructurePreservation:
    def test_xlsx_preserves_columns_with_empty_cells(self):
        # Middle column empty: column alignment (and column letters) must stay intact so that
        # evidence like "Spalte C" remains meaningful.
        content = _make_xlsx_bytes({"VVT": [["Zweck", None, "Rechtsgrundlage"], ["Lohn", None, "Art. 6"]]})
        text = extract_text_from_xlsx(content)
        # Column-letter header row present (with the leading row-number column)
        assert "| Zeile | A | B | C |" in text
        # Rows carry their 1-based number and keep the empty middle column.
        assert "| 1 | Zweck |  | Rechtsgrundlage |" in text
        assert "| 2 | Lohn |  | Art. 6 |" in text

    def test_xlsx_rows_are_numbered_sequentially(self):
        content = _make_xlsx_bytes({"Daten": [["a"], ["b"], ["c"]]})
        text = extract_text_from_xlsx(content)
        assert "| 1 | a |" in text
        assert "| 2 | b |" in text
        assert "| 3 | c |" in text

    def test_docx_table_rendered_as_markdown(self):
        content = _make_docx_bytes([], table_rows=[["Feld", "Wert"], ["Zweck", "Lohnabrechnung"]])
        text = extract_text_from_docx(content)
        assert "| Feld | Wert |" in text
        assert "| --- | --- |" in text
        assert "| Zweck | Lohnabrechnung |" in text


def _make_pdf_bytes(pages: list[list[str]]) -> bytes:
    """Build a digital (text-layer) PDF with one entry per page; empty list = blank page."""
    import fitz

    doc = fitz.open()
    for lines in pages:
        page = doc.new_page()
        y = 72.0
        for line in lines:
            page.insert_text((72, y), line, fontsize=11)
            y += 18
    data = doc.tobytes()
    doc.close()
    return data


# Enough text per page to stay above ocr_min_chars_per_page (default 50) → no OCR triggered.
_PAGE1_LINES = ["Datenschutzhinweise zur Lohnabrechnung im Sinne der DSGVO," , "Verantwortlicher ist die Muster GmbH mit Sitz in Frankfurt am Main."]
_PAGE2_LINES = ["Speicherdauer: zehn Jahre gemaess HGB und Abgabenordnung.", "Eine Uebermittlung in Drittlaender findet nicht statt, Stand 2026."]


class TestPdfPageAnchors:
    def test_multi_page_pdf_carries_page_anchors(self):
        result = extract_text("hinweis.pdf", _make_pdf_bytes([_PAGE1_LINES, _PAGE2_LINES]))
        assert "[Seite 1]" in result.text
        assert "[Seite 2]" in result.text
        # Content sits under the right anchor.
        assert result.text.index("[Seite 1]") < result.text.index("Muster GmbH")
        assert result.text.index("Muster GmbH") < result.text.index("[Seite 2]")
        assert result.text.index("[Seite 2]") < result.text.index("zehn Jahre")
        assert result.page_count == 2

    def test_single_page_pdf_has_no_anchor(self):
        result = extract_text("hinweis.pdf", _make_pdf_bytes([_PAGE1_LINES]))
        assert "[Seite" not in result.text
        assert "Muster GmbH" in result.text

    def test_ocr_merged_pages_keep_anchors(self, monkeypatch):
        import app.services.document_processor as dp

        monkeypatch.setattr(dp.settings, "ollama_ocr_enabled", True, raising=False)
        monkeypatch.setattr(dp.settings, "ocr_per_page", True, raising=False)
        ocr_text = "OCR-Inhalt der gescannten Seite zwei mit ausreichend vielen Zeichen fuer die Schwelle."
        monkeypatch.setattr(dp, "_ocr_pages", lambda content, targets: dict.fromkeys(targets, ocr_text))

        # Page 1 digital, page 2 blank (sparse → OCR).
        result = extract_text("scan.pdf", _make_pdf_bytes([_PAGE1_LINES, []]))
        assert result.extraction_method == "ocr"
        assert "[Seite 1]" in result.text and "[Seite 2]" in result.text
        assert result.text.index("[Seite 2]") < result.text.index("OCR-Inhalt")

    def test_assemble_pages_skips_empty_but_keeps_real_page_numbers(self):
        import app.services.document_processor as dp

        text = dp._assemble_pages(["Erste Seite", "", "Dritte Seite"], 3)
        assert "[Seite 1]" in text
        assert "[Seite 2]" not in text  # empty page emits no anchor
        assert "[Seite 3]" in text  # numbering reflects the real page index

    def test_assemble_pages_single_page_unanchored(self):
        import app.services.document_processor as dp

        assert dp._assemble_pages(["Inhalt"], 1) == "Inhalt"


class TestLegacyDoc:
    def test_doc_without_converter_raises(self, monkeypatch):
        import app.services.document_processor as dp

        # Simulate no converter installed → clear error instead of silent empty text.
        monkeypatch.setattr(dp.shutil, "which", lambda _name: None)
        with pytest.raises(dp.UnsupportedDocumentError):
            extract_text("altes_dokument.doc", b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1fakeole2")


class _FakeResp:
    """OpenAI-compatible chat-completions response shape (works for Ollama /v1, vLLM, llama.cpp)."""

    def __init__(self, content: str):
        self._content = content

    def raise_for_status(self):
        return None

    def json(self):
        return {"choices": [{"message": {"content": self._content}}]}


class TestOcrRetry:
    """A3: per-page OCR is retried on transient failure / empty output instead of silently lost."""

    def _patch_render_and_sleep(self, monkeypatch):
        monkeypatch.setattr(dp, "_pdf_page_to_png_bytes", lambda content, i, dpi=150: b"x")
        monkeypatch.setattr(dp.time, "sleep", lambda _s: None)

    def test_retries_until_attempts_exhausted_then_empty(self, monkeypatch):
        self._patch_render_and_sleep(monkeypatch)
        monkeypatch.setattr(dp.settings, "ocr_retry_attempts", 3, raising=False)
        calls = {"n": 0}

        def _always_fail(*a, **k):
            calls["n"] += 1
            raise RuntimeError("boom")

        monkeypatch.setattr(dp.httpx, "post", _always_fail)
        result = dp._ocr_pages(b"fakepdf", [0])
        assert result == {0: ""}
        assert calls["n"] == 3  # one attempt per configured retry

    def test_succeeds_after_a_transient_failure(self, monkeypatch):
        self._patch_render_and_sleep(monkeypatch)
        monkeypatch.setattr(dp.settings, "ocr_retry_attempts", 2, raising=False)
        calls = {"n": 0}

        def _fail_then_ok(*a, **k):
            calls["n"] += 1
            if calls["n"] == 1:
                raise RuntimeError("transient")
            return _FakeResp("Seiteninhalt erkannt")

        monkeypatch.setattr(dp.httpx, "post", _fail_then_ok)
        result = dp._ocr_pages(b"fakepdf", [0])
        assert result == {0: "Seiteninhalt erkannt"}
        assert calls["n"] == 2

    def test_empty_response_is_retried(self, monkeypatch):
        self._patch_render_and_sleep(monkeypatch)
        monkeypatch.setattr(dp.settings, "ocr_retry_attempts", 2, raising=False)
        calls = {"n": 0}

        def _empty_then_ok(*a, **k):
            calls["n"] += 1
            return _FakeResp("" if calls["n"] == 1 else "Treffer")

        monkeypatch.setattr(dp.httpx, "post", _empty_then_ok)
        result = dp._ocr_pages(b"fakepdf", [0])
        assert result == {0: "Treffer"}
        assert calls["n"] == 2


class TestOcrEndpointResolution:
    """OCR runs against the OpenAI-compatible chat-completions API of a configurable backend."""

    def test_default_is_ollama_v1(self, monkeypatch):
        monkeypatch.setattr(dp.settings, "ocr_base_url", "", raising=False)
        monkeypatch.setattr(dp.settings, "llm_provider", "ollama", raising=False)
        monkeypatch.setattr(dp.settings, "ollama_base_url", "http://localhost:11434", raising=False)
        url, headers = dp._ocr_chat_endpoint()
        assert url == "http://localhost:11434/v1/chat/completions"
        assert headers == {}

    def test_explicit_override_with_api_key(self, monkeypatch):
        from pydantic import SecretStr

        monkeypatch.setattr(dp.settings, "ocr_base_url", "http://vision:8000/v1/", raising=False)
        monkeypatch.setattr(dp.settings, "ocr_api_key", SecretStr("s3cret"), raising=False)
        url, headers = dp._ocr_chat_endpoint()
        assert url == "http://vision:8000/v1/chat/completions"
        assert headers == {"Authorization": "Bearer s3cret"}

    def test_openai_compatible_provider_is_used_when_no_override(self, monkeypatch):
        from pydantic import SecretStr

        monkeypatch.setattr(dp.settings, "ocr_base_url", "", raising=False)
        monkeypatch.setattr(dp.settings, "ocr_api_key", SecretStr(""), raising=False)
        monkeypatch.setattr(dp.settings, "llm_provider", "openai_compatible", raising=False)
        monkeypatch.setattr(dp.settings, "llm_base_url", "http://vllm:8000", raising=False)
        monkeypatch.setattr(dp.settings, "llm_api_key", SecretStr("vllm-key"), raising=False)
        url, headers = dp._ocr_chat_endpoint()
        assert url == "http://vllm:8000/v1/chat/completions"
        assert headers == {"Authorization": "Bearer vllm-key"}

    def test_ocr_model_falls_back_to_legacy_setting(self, monkeypatch):
        monkeypatch.setattr(dp.settings, "ocr_model", "", raising=False)
        monkeypatch.setattr(dp.settings, "ollama_ocr_model", "qwen2.5-vl", raising=False)
        assert dp._ocr_model_name() == "qwen2.5-vl"
        monkeypatch.setattr(dp.settings, "ocr_model", "Qwen/Qwen2.5-VL-7B-Instruct", raising=False)
        assert dp._ocr_model_name() == "Qwen/Qwen2.5-VL-7B-Instruct"

    def test_payload_uses_openai_multimodal_format(self, monkeypatch):
        monkeypatch.setattr(dp, "_pdf_page_to_png_bytes", lambda content, i, dpi=150: b"img")
        monkeypatch.setattr(dp.time, "sleep", lambda _s: None)
        monkeypatch.setattr(dp.settings, "ocr_retry_attempts", 1, raising=False)
        captured = {}

        def _capture(url, json=None, headers=None, timeout=None):
            captured["url"] = url
            captured["payload"] = json
            return _FakeResp("Erkannter Inhalt der gescannten Seite")

        monkeypatch.setattr(dp.httpx, "post", _capture)
        result = dp._ocr_pages(b"fakepdf", [0])
        assert result == {0: "Erkannter Inhalt der gescannten Seite"}
        msg = captured["payload"]["messages"][0]
        parts = msg["content"]
        assert {p["type"] for p in parts} == {"text", "image_url"}
        image_part = next(p for p in parts if p["type"] == "image_url")
        assert image_part["image_url"]["url"].startswith("data:image/png;base64,")
        assert captured["payload"]["temperature"] == 0
        assert captured["url"].endswith("/chat/completions")


class TestOcrQualitySignal:
    def test_digital_formats_report_zero_low_quality_pages(self):
        # Non-paged / digital extraction carries the default 0 (no OCR pages, nothing lost).
        result = extract_text("report.docx", _make_docx_bytes(["Inhalt"]))
        assert result.ocr_low_quality_pages == 0


def _make_png_bytes() -> bytes:
    """Build a tiny in-memory PNG via PyMuPDF (no Pillow dependency)."""
    import fitz

    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 12, 12), False)
    pix.clear_with(255)
    return pix.tobytes("png")


class TestImageUpload:
    """A4: standalone scanned images route through the PDF/OCR path."""

    def test_png_routes_through_pdf_path_without_crashing(self, monkeypatch):
        # OCR disabled → empty text, but conversion + routing must succeed (no exception).
        monkeypatch.setattr(dp.settings, "ollama_ocr_enabled", False, raising=False)
        result = extract_text("scan.png", _make_png_bytes())
        assert isinstance(result, ExtractionResult)
        assert isinstance(result.text, str)

    def test_unknown_image_bytes_do_not_raise(self, monkeypatch):
        monkeypatch.setattr(dp.settings, "ollama_ocr_enabled", False, raising=False)
        result = extract_text("scan.png", b"not a real image")
        assert isinstance(result, ExtractionResult)


def _make_encrypted_pdf_bytes() -> bytes:
    """Build a password-protected (AES-256) PDF that requires a user password to read."""
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Vertraulicher Inhalt")
    out = doc.tobytes(
        encryption=fitz.PDF_ENCRYPT_AES_256,
        user_pw="secret",
        owner_pw="owner",
    )
    doc.close()
    return out


def _make_image_heavy_pdf_bytes() -> bytes:
    """Build a PDF page that is mostly a scanned-style image with only a tiny header text."""
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 400, 520), False)
    pix.clear_with(220)
    page.insert_image(page.rect, pixmap=pix)
    page.insert_text((72, 18), "Kopfzeile Az. DS-1")  # short digital text → not 'sparse' by chars
    out = doc.tobytes()
    doc.close()
    return out


class TestEncryptedPdf:
    """A3: password-protected PDFs raise a clear error instead of yielding empty text."""

    def test_encrypted_pdf_raises(self):
        with pytest.raises(dp.UnsupportedDocumentError):
            dp.extract_text_from_pdf(_make_encrypted_pdf_bytes())


class TestImageHeavyPages:
    """A2: pages dominated by an image but with a little digital text are OCR candidates."""

    def test_image_heavy_page_detected(self):
        import fitz

        with fitz.open(stream=_make_image_heavy_pdf_bytes(), filetype="pdf") as doc:
            assert dp._image_heavy_pages(doc) == {0}

    def test_text_page_not_flagged(self):
        import fitz

        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Ausführlicher Fließtext " * 40)  # plenty of digital text
        data = doc.tobytes()
        doc.close()
        with fitz.open(stream=data, filetype="pdf") as reopened:
            assert dp._image_heavy_pages(reopened) == set()

    def test_disabled_via_threshold(self, monkeypatch):
        import fitz

        monkeypatch.setattr(dp.settings, "ocr_image_area_threshold", 1.0, raising=False)
        with fitz.open(stream=_make_image_heavy_pdf_bytes(), filetype="pdf") as doc:
            assert dp._image_heavy_pages(doc) == set()


class TestDetectLanguage:
    """A4: lightweight DE/EN language detection for the LLM language hint."""

    def test_german_text(self):
        text = (
            "Die Verarbeitung der personenbezogenen Daten erfolgt zur Lohnabrechnung und wird "
            "auf Grundlage der gesetzlichen Vorgaben durchgeführt, damit die Fristen eingehalten werden."
        )
        assert dp.detect_language(text) == "de"

    def test_english_text(self):
        text = (
            "The processing of the personal data is carried out for payroll purposes and shall "
            "be based on the legal requirements so that the retention periods are not exceeded."
        )
        assert dp.detect_language(text) == "en"

    def test_too_short_returns_none(self):
        assert dp.detect_language("Lohn") is None
        assert dp.detect_language("") is None
        assert dp.detect_language(None) is None


class TestImageMagicBytes:
    """A4: magic-byte validation accepts real image headers and rejects mismatches."""

    def _verify(self):
        from app.api.routes.documents import _verify_magic_bytes
        return _verify_magic_bytes

    def test_accepts_matching_image_headers(self):
        from fastapi import HTTPException

        verify = self._verify()
        verify(b"\xff\xd8\xff\xe0\x00\x10JFIF", "jpg", "s.jpg")
        verify(b"\x89PNG\r\n\x1a\n\x00\x00", "png", "s.png")
        verify(b"II*\x00\x08\x00\x00\x00", "tiff", "s.tiff")
        verify(b"MM\x00*\x00\x00\x00\x08", "tiff", "s.tiff")

    def test_rejects_mismatched_header(self):
        from fastapi import HTTPException

        verify = self._verify()
        with pytest.raises(HTTPException) as exc:
            verify(b"%PDF-1.7 ...", "png", "fake.png")
        assert exc.value.status_code == 415


def _make_pptx_bytes(slides: list[dict]) -> bytes:
    """Build an in-memory PPTX. Each slide dict: {"text": str|None, "table": rows|None, "notes": str|None}."""
    import pptx
    from pptx.util import Inches

    prs = pptx.Presentation()
    blank = prs.slide_layouts[6]
    for spec in slides:
        slide = prs.slides.add_slide(blank)
        if spec.get("text"):
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            box.text_frame.text = spec["text"]
        if spec.get("table"):
            rows = spec["table"]
            table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(1), Inches(3), Inches(8), Inches(2)).table
            for r, row in enumerate(rows):
                for c, val in enumerate(row):
                    table.cell(r, c).text = val
        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestExtractTextFromPptx:
    def test_multi_slide_anchors_with_empty_first_slide(self):
        from app.services.document_processor import extract_text_from_pptx

        content = _make_pptx_bytes([{}, {"text": "Löschkonzept nach sechs Monaten"}])
        text = extract_text_from_pptx(content)
        assert "[Folie 1]" not in text  # empty slide emits no anchor …
        assert "[Folie 2]" in text      # … but numbering reflects the real slide index
        assert "Löschkonzept" in text

    def test_single_slide_has_no_anchor(self):
        from app.services.document_processor import extract_text_from_pptx

        text = extract_text_from_pptx(_make_pptx_bytes([{"text": "TOM-Übersicht"}]))
        assert "[Folie" not in text
        assert "TOM-Übersicht" in text

    def test_notes_and_table_extracted(self):
        from app.services.document_processor import extract_text_from_pptx

        content = _make_pptx_bytes([
            {"table": [["Maßnahme", "Status"], ["Pseudo|nymisierung", "umgesetzt"]], "notes": "AVV liegt vor."},
            {"text": "Zweite Folie"},
        ])
        text = extract_text_from_pptx(content)
        assert "| Maßnahme | Status |" in text
        assert "Pseudo\\|nymisierung" in text  # pipe escaped like DOCX/XLSX cells
        assert "--- Notizen ---" in text
        assert "AVV liegt vor." in text

    def test_dispatcher_routes_pptx(self):
        result = extract_text("folien.pptx", _make_pptx_bytes([{"text": "Inhalt der Präsentation"}]))
        assert result.extraction_method == EXTRACTION_METHOD_TEXT
        assert "Inhalt der Präsentation" in result.text

    def test_corrupt_pptx_raises(self):
        import zipfile

        with pytest.raises(zipfile.BadZipFile):
            extract_text("kaputt.pptx", b"PK\x03\x04 not a real pptx")


class TestExtractTextFromCsv:
    def test_semicolon_delimiter_sniffed(self):
        from app.services.document_processor import extract_text_from_csv

        text = extract_text_from_csv("Zweck;Rechtsgrundlage\nLohn;Art. 6\n".encode())
        assert "| Zeile | A | B |" in text
        assert "| 1 | Zweck | Rechtsgrundlage |" in text
        assert "| 2 | Lohn | Art. 6 |" in text

    def test_comma_and_tab_delimiters(self):
        from app.services.document_processor import extract_text_from_csv

        comma = extract_text_from_csv(b"a,b\nc,d\n")
        assert "| 1 | a | b |" in comma
        tab = extract_text_from_csv(b"a\tb\nc\td\n")
        assert "| 1 | a | b |" in tab

    def test_cp1252_umlauts_and_bom(self):
        from app.services.document_processor import extract_text_from_csv

        cp = extract_text_from_csv("Tätigkeit;Dauer\nLöschung;6 Monate\n".encode("cp1252"))
        assert "Tätigkeit" in cp and "Löschung" in cp
        bom = extract_text_from_csv("﻿Spalte1;Spalte2\n".encode("utf-8"))
        assert "| 1 | Spalte1 | Spalte2 |" in bom  # BOM must not stick to the first cell

    def test_ragged_rows_padded_and_pipes_escaped(self):
        from app.services.document_processor import extract_text_from_csv

        text = extract_text_from_csv(b"a;b;c\nx\nmit|pipe;y;z\n")
        assert "| 2 | x |  |  |" in text
        assert "mit\\|pipe" in text

    def test_empty_csv_returns_empty(self):
        from app.services.document_processor import extract_text_from_csv

        assert extract_text_from_csv(b"") == ""
        assert extract_text_from_csv(b"   \n  ") == ""

    def test_dispatcher_routes_csv(self):
        result = extract_text("export.csv", b"Zweck;Dauer\nLohn;10 Jahre\n")
        assert result.extraction_method == EXTRACTION_METHOD_TEXT
        assert "| Zeile | A | B |" in result.text
