"""Text extraction from documents (PDF, DOCX, XLSX, PPTX, CSV).

PDF and image extraction is delegated to pdf_extractor (which in turn uses
ocr_service for vision-LLM OCR). This module retains the dispatcher, language
detection, and the DOCX/XLSX/PPTX/CSV/DOC extractors.
"""

import csv
import io
import logging
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

import docx
import openpyxl
from docx.oxml.ns import qn
from lxml import etree
from openpyxl.utils import get_column_letter

from app.services.pdf_extractor import (
    EXTRACTION_METHOD_OCR,  # noqa: F401  (re-exported for callers)
    EXTRACTION_METHOD_TEXT,
    ExtractionResult,
    UnsupportedDocumentError,
    _assemble_anchored,
    extract_text_from_image,
    extract_text_from_pdf,
)

logger = logging.getLogger(__name__)

# Re-export so existing callers that import these from document_processor still work.
__all__ = [
    "ExtractionResult",
    "UnsupportedDocumentError",
    "EXTRACTION_METHOD_TEXT",
    "EXTRACTION_METHOD_OCR",
    "detect_language",
    "extract_text",
    "extract_text_from_pdf",
    "extract_text_from_image",
    "extract_text_from_docx",
    "extract_text_from_xlsx",
    "extract_text_from_pptx",
    "extract_text_from_csv",
    "extract_text_from_doc",
]

# ---------------------------------------------------------------------------
# Language detection
# ---------------------------------------------------------------------------

_DE_STOPWORDS = frozenset(
    {
        "der",
        "die",
        "das",
        "und",
        "ist",
        "nicht",
        "werden",
        "wird",
        "für",
        "mit",
        "auch",
        "sich",
        "eine",
        "einer",
        "von",
        "den",
        "dem",
        "des",
        "im",
        "zur",
        "zum",
        "bei",
        "durch",
    }
)
_EN_STOPWORDS = frozenset(
    {
        "the",
        "and",
        "are",
        "with",
        "for",
        "this",
        "that",
        "from",
        "have",
        "will",
        "not",
        "shall",
        "been",
        "their",
        "which",
        "such",
        "any",
        "into",
        "data",
        "processing",
    }
)


def detect_language(text: str | None) -> str | None:
    """Best-effort DE/EN detection via stop-word frequency. Returns "de", "en", or None."""
    if not text:
        return None
    words = re.findall(r"[a-zäöüß]+", text.lower())
    if len(words) < 20:
        return None
    de = sum(1 for w in words if w in _DE_STOPWORDS)
    en = sum(1 for w in words if w in _EN_STOPWORDS)
    if de == 0 and en == 0:
        return None
    if de >= en * 1.2:
        return "de"
    if en >= de * 1.2:
        return "en"
    return None


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------


def _rows_to_markdown(rows: list[list[str]]) -> str:
    """Render pre-extracted cell rows as a Markdown table (first row = header)."""
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * width) + " |"
    body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join([header, sep, *body])


def _escape_cell(value: str | None) -> str:
    """Escape a table cell for Markdown output (collapse newlines, escape pipes)."""
    return (value or "").replace("\n", " ").replace("|", "\\|").strip()


def _docx_table_to_markdown(table) -> str:
    rows = [[_escape_cell(cell.text) for cell in row.cells] for row in table.rows]
    return _rows_to_markdown(rows)


def _iter_header_footer_text(doc) -> list[str]:
    """Collect text from every section header/footer."""
    parts: list[str] = []
    seen: set[str] = set()
    for section in doc.sections:
        for hf in (
            section.header,
            section.first_page_header,
            section.even_page_header,
            section.footer,
            section.first_page_footer,
            section.even_page_footer,
        ):
            if hf is None:
                continue
            try:
                for para in hf.paragraphs:
                    txt = (para.text or "").strip()
                    if txt and txt not in seen:
                        seen.add(txt)
                        parts.append(txt)
                for table in hf.tables:
                    md = _docx_table_to_markdown(table)
                    if md and md not in seen:
                        seen.add(md)
                        parts.append(md)
            except Exception as exc:
                logger.debug("DOCX header/footer extraction skipped a part: %s", exc)
    return parts


def _iter_textbox_text(doc) -> list[str]:
    """Collect text inside drawing text boxes (``w:txbxContent``)."""
    parts: list[str] = []
    try:
        for txbx in doc.element.body.iter(qn("w:txbxContent")):
            for p in txbx.iter(qn("w:p")):
                txt = "".join((t.text or "") for t in p.iter(qn("w:t"))).strip()
                if txt:
                    parts.append(txt)
    except Exception as exc:
        logger.debug("DOCX text-box extraction skipped: %s", exc)
    return parts


def _iter_footnote_text(doc) -> list[str]:
    """Collect footnote/endnote text from the related XML parts."""
    parts: list[str] = []
    try:
        rels = doc.part.rels.values()
    except Exception:
        return parts
    for rel in rels:
        rtype = getattr(rel, "reltype", "") or ""
        if not (rtype.endswith("/footnotes") or rtype.endswith("/endnotes")):
            continue
        if getattr(rel, "is_external", False):
            continue
        try:
            root = etree.fromstring(rel.target_part.blob)
        except Exception as exc:
            logger.debug("DOCX footnote part unreadable: %s", exc)
            continue
        note_tag = qn("w:footnote") if rtype.endswith("/footnotes") else qn("w:endnote")
        for note in root.iter(note_tag):
            if note.get(qn("w:type")) in ("separator", "continuationSeparator"):
                continue
            txt = "".join((t.text or "") for t in note.iter(qn("w:t"))).strip()
            if txt:
                parts.append(txt)
    return parts


def extract_text_from_docx(content: bytes) -> str:
    """Extract text from DOCX bytes including text boxes, headers/footers and footnotes."""
    doc = docx.Document(io.BytesIO(content))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text and para.text.strip():
            parts.append(para.text)
    parts.extend(_iter_textbox_text(doc))
    for table in doc.tables:
        md = _docx_table_to_markdown(table)
        if md:
            parts.append(md)
    header_footer = _iter_header_footer_text(doc)
    if header_footer:
        parts.append("--- Kopf-/Fußzeilen ---")
        parts.extend(header_footer)
    footnotes = _iter_footnote_text(doc)
    if footnotes:
        parts.append("--- Fußnoten ---")
        parts.extend(footnotes)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# XLSX extraction
# ---------------------------------------------------------------------------


def extract_text_from_xlsx(content: bytes) -> str:
    """Extract text from XLSX bytes as one Markdown table per sheet."""
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    out: list[str] = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        out.append(f"--- Sheet: {sheet} ---")
        rows = list(ws.iter_rows(values_only=True))
        max_cols = max((len(r) for r in rows), default=0)
        if max_cols == 0:
            continue
        col_header = ["Zeile"] + [get_column_letter(c + 1) for c in range(max_cols)]
        out.append("| " + " | ".join(col_header) + " |")
        out.append("| " + " | ".join(["---"] * (max_cols + 1)) + " |")
        for row_no, r in enumerate(rows, start=1):
            cells = [
                "" if c is None else str(c).replace("\n", " ").replace("|", "\\|")
                for c in r
            ]
            cells += [""] * (max_cols - len(cells))
            out.append("| " + " | ".join([str(row_no), *cells]) + " |")
    return "\n".join(out)


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------


def _slide_anchor(slide_index: int) -> str:
    """Anchor line marking the start of a PPTX slide in the extracted text (1-based)."""
    return f"[Folie {slide_index + 1}]"


def _pptx_shape_texts(shapes, depth: int = 0) -> list[str]:
    """Collect text frames and tables from PPTX shapes, recursing one level into groups."""
    parts: list[str] = []
    for shape in shapes:
        try:
            if getattr(shape, "has_text_frame", False):
                txt = (shape.text_frame.text or "").strip()
                if txt:
                    parts.append(txt)
            if getattr(shape, "has_table", False):
                rows = [
                    [_escape_cell(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                md = _rows_to_markdown(rows)
                if md:
                    parts.append(md)
            if depth < 1:
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    parts.extend(_pptx_shape_texts(shape.shapes, depth + 1))
        except Exception as exc:
            logger.debug("PPTX shape extraction skipped a shape: %s", exc)
    return parts


def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PPTX bytes: slide text frames, tables and speaker notes."""
    import pptx

    prs = pptx.Presentation(io.BytesIO(content))
    slides = list(prs.slides)
    slide_parts: list[str] = []
    for slide in slides:
        parts = _pptx_shape_texts(slide.shapes)
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append("--- Notizen ---")
                    parts.append(notes)
        except Exception as exc:
            logger.debug("PPTX notes extraction skipped: %s", exc)
        slide_parts.append("\n".join(parts))
    return _assemble_anchored(slide_parts, len(slides), _slide_anchor)


# ---------------------------------------------------------------------------
# CSV extraction
# ---------------------------------------------------------------------------


def _sniff_csv_delimiter(sample: str) -> str:
    """Detect the CSV delimiter (German exports are predominantly ';')."""
    try:
        return csv.Sniffer().sniff(sample, delimiters=",;\t|").delimiter
    except csv.Error:
        first_line = sample.splitlines()[0] if sample.splitlines() else ""
        counts = {d: first_line.count(d) for d in (";", ",", "\t", "|")}
        best = max(counts, key=lambda d: counts[d])
        return best if counts[best] > 0 else ","


def extract_text_from_csv(content: bytes) -> str:
    """Extract CSV bytes as a Markdown table in the same shape as XLSX sheets."""
    text = _decode_text_bytes(content).lstrip("﻿")
    if not text.strip():
        return ""
    delimiter = _sniff_csv_delimiter(text[:4096])
    rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
    max_cols = max((len(r) for r in rows), default=0)
    if max_cols == 0:
        return ""
    out: list[str] = []
    col_header = ["Zeile"] + [get_column_letter(c + 1) for c in range(max_cols)]
    out.append("| " + " | ".join(col_header) + " |")
    out.append("| " + " | ".join(["---"] * (max_cols + 1)) + " |")
    for row_no, r in enumerate(rows, start=1):
        cells = [_escape_cell(c) for c in r]
        cells += [""] * (max_cols - len(cells))
        out.append("| " + " | ".join([str(row_no), *cells]) + " |")
    return "\n".join(out)


# ---------------------------------------------------------------------------
# Legacy .doc and text fallback
# ---------------------------------------------------------------------------


def extract_text_from_doc(content: bytes) -> str:
    """Extract text from a legacy .doc (OLE2) file via LibreOffice or antiword."""
    if shutil.which("antiword"):
        try:
            return subprocess.run(
                ["antiword", "-"], input=content, capture_output=True, timeout=120
            ).stdout.decode("utf-8", errors="replace")
        except Exception as exc:
            logger.warning("antiword .doc extraction failed: %s", exc)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        try:
            with tempfile.TemporaryDirectory() as tmp:
                src = Path(tmp) / "in.doc"
                src.write_bytes(content)
                subprocess.run(
                    [
                        soffice,
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        tmp,
                        str(src),
                    ],
                    capture_output=True,
                    timeout=180,
                )
                out = Path(tmp) / "in.docx"
                if out.exists():
                    return extract_text_from_docx(out.read_bytes())
        except Exception as exc:
            logger.warning("LibreOffice .doc conversion failed: %s", exc)
    raise UnsupportedDocumentError(
        "Legacy .doc files require a converter (LibreOffice/soffice or antiword) which is not "
        "available. Please upload the document as .docx or PDF."
    )


def _decode_text_bytes(content: bytes) -> str:
    """Decode arbitrary text bytes: UTF-8 → charset detection → CP1252."""
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        pass
    if b"\x00" in content[:65536]:
        return ""
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(content).best()
        if best is not None:
            return str(best)
    except Exception:
        pass
    try:
        return content.decode("cp1252")
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_IMAGE_EXT_TO_FITZ = {
    "jpg": "jpg",
    "jpeg": "jpg",
    "png": "png",
    "tif": "tiff",
    "tiff": "tiff",
}


def extract_text(filename: str, content: bytes) -> ExtractionResult:
    """Dispatcher for text extraction. Returns ExtractionResult."""
    filename_lower = filename.lower()
    if filename_lower.endswith(".pdf"):
        return extract_text_from_pdf(content)
    ext = filename_lower.rsplit(".", 1)[-1] if "." in filename_lower else ""
    if ext in _IMAGE_EXT_TO_FITZ:
        try:
            return extract_text_from_image(content, _IMAGE_EXT_TO_FITZ[ext])
        except Exception:
            logger.exception(
                "Image text extraction failed", extra={"doc_filename": filename}
            )
            raise
    if filename_lower.endswith(".docx"):
        try:
            text = extract_text_from_docx(content)
            return ExtractionResult(
                text=text,
                extraction_method=EXTRACTION_METHOD_TEXT,
                char_count=len(text),
            )
        except Exception:
            logger.exception(
                "DOCX text extraction failed", extra={"doc_filename": filename}
            )
            raise
    if filename_lower.endswith(".xlsx"):
        try:
            text = extract_text_from_xlsx(content)
            return ExtractionResult(
                text=text,
                extraction_method=EXTRACTION_METHOD_TEXT,
                char_count=len(text),
            )
        except Exception:
            logger.exception(
                "XLSX text extraction failed", extra={"doc_filename": filename}
            )
            raise
    if filename_lower.endswith(".pptx"):
        try:
            text = extract_text_from_pptx(content)
            return ExtractionResult(
                text=text,
                extraction_method=EXTRACTION_METHOD_TEXT,
                char_count=len(text),
            )
        except Exception:
            logger.exception(
                "PPTX text extraction failed", extra={"doc_filename": filename}
            )
            raise
    if filename_lower.endswith(".csv"):
        try:
            text = extract_text_from_csv(content)
            return ExtractionResult(
                text=text,
                extraction_method=EXTRACTION_METHOD_TEXT,
                char_count=len(text),
            )
        except Exception:
            logger.exception(
                "CSV text extraction failed", extra={"doc_filename": filename}
            )
            raise
    if filename_lower.endswith(".doc"):
        text = extract_text_from_doc(content)
        return ExtractionResult(
            text=text, extraction_method=EXTRACTION_METHOD_TEXT, char_count=len(text)
        )
    text = _decode_text_bytes(content)
    if not text:
        logger.debug("No decodable text for %s", filename)
    return ExtractionResult(
        text=text, extraction_method=EXTRACTION_METHOD_TEXT, char_count=len(text)
    )
