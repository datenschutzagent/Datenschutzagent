"""In-memory gold sample documents + expectations for the extraction evals.

Documents are generated programmatically (no binary fixtures in the repo) so the gold set is
transparent and diffable. Each sample declares the tokens that must survive extraction and,
for tabular formats, a column-header signature that must be preserved.
"""
from __future__ import annotations

import io
from collections.abc import Callable
from dataclasses import dataclass


@dataclass(frozen=True)
class ExtractionSample:
    key: str
    filename: str
    builder: Callable[[], bytes]
    expected_tokens: list[str]   # must all appear in the extracted text
    column_header: str | None = None  # e.g. "| A | B | C |" for tabular formats
    min_chars: int = 1


def _digital_pdf_bytes() -> bytes:
    """A digital-born (text-layer) PDF — exercises the offline PDF→Markdown extraction path.

    Built with PyMuPDF so the page carries a real text layer (no OCR needed). This closes the gap
    that the gold set previously only covered DOCX/XLSX while the PDF path — the most common upload
    format — went unmeasured.
    """
    import fitz  # PyMuPDF

    lines = [
        "Datenschutzhinweise zur Lohnabrechnung",
        "Verantwortlicher: Muster GmbH",
        "Zweck der Verarbeitung: Lohn- und Gehaltsabrechnung",
        "Rechtsgrundlage: Art. 6 Abs. 1 lit. c DSGVO",
        "Speicherdauer: 10 Jahre gemaess HGB und AO",
        "Eine Uebermittlung in Drittlaender findet nicht statt.",
    ]
    doc = fitz.open()
    page = doc.new_page()
    y = 72.0
    for line in lines:
        page.insert_text((72, y), line, fontsize=11)
        y += 24
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def _vvt_xlsx_bytes() -> bytes:
    """A small VVT/ROPA sheet with an intentionally empty middle column."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "VVT"
    rows = [
        ["Verarbeitungstätigkeit", None, "Rechtsgrundlage", "Speicherdauer"],
        ["Lohnabrechnung", None, "Art. 6 Abs. 1 lit. c DSGVO", "10 Jahre"],
        ["Bewerbermanagement", None, "Art. 6 Abs. 1 lit. b DSGVO", "6 Monate"],
    ]
    for r in rows:
        ws.append(r)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _avv_docx_bytes() -> bytes:
    """A small AVV-style DOCX containing a paragraph and a table."""
    import docx

    doc = docx.Document()
    doc.add_paragraph("Auftragsverarbeitungsvertrag gemäß Art. 28 DSGVO")
    doc.add_paragraph("Der Auftragsverarbeiter verarbeitet personenbezogene Daten weisungsgebunden.")
    table = doc.add_table(rows=3, cols=2)
    cells = [["Gegenstand", "Lohnbuchhaltung"], ["Subunternehmer", "Keine"], ["Drittland", "Nein"]]
    for r_idx, row in enumerate(cells):
        for c_idx, val in enumerate(row):
            table.rows[r_idx].cells[c_idx].text = val
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _tom_pptx_bytes() -> bytes:
    """A two-slide TOM deck: slide 1 with a table, slide 2 with body text and speaker notes."""
    import pptx
    from pptx.util import Inches

    prs = pptx.Presentation()
    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    rows, cols = 3, 2
    table = slide1.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(8), Inches(3)).table
    cells = [["Maßnahme", "Status"], ["Verschlüsselung", "umgesetzt"], ["Zugriffskontrolle", "geplant"]]
    for r, row in enumerate(cells):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    slide2 = prs.slides.add_slide(blank)
    box = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text_frame.text = "Löschkonzept: Daten werden nach 6 Monaten geloescht."
    slide2.notes_slide.notes_text_frame.text = "Hinweis: AVV mit Hosting-Anbieter liegt vor."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _vvt_csv_bytes() -> bytes:
    """A semicolon-delimited German CSV export (umlauts, typical Excel style)."""
    import csv

    out = io.StringIO()
    writer = csv.writer(out, delimiter=";")
    writer.writerow(["Verarbeitungstätigkeit", "Rechtsgrundlage", "Speicherdauer"])
    writer.writerow(["Lohnabrechnung", "Art. 6 Abs. 1 lit. c DSGVO", "10 Jahre"])
    writer.writerow(["Bewerbermanagement", "Art. 6 Abs. 1 lit. b DSGVO", "6 Monate"])
    return out.getvalue().encode("utf-8")


def _header_footer_docx_bytes() -> bytes:
    """A DOCX whose controller / version metadata lives only in the header and footer."""
    import docx

    doc = docx.Document()
    doc.add_paragraph("Verarbeitungstätigkeit: Lohnabrechnung")
    section = doc.sections[0]
    section.header.paragraphs[0].text = "Verantwortlicher: Muster GmbH"
    section.footer.paragraphs[0].text = "Stand: 01.01.2026 — Az. DS-2026-042"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


SAMPLES: dict[str, ExtractionSample] = {
    "digital_pdf": ExtractionSample(
        key="digital_pdf",
        filename="lohnabrechnung.pdf",
        builder=_digital_pdf_bytes,
        # Controller, purpose, legal basis and retention period must all survive PDF extraction.
        expected_tokens=[
            "Muster GmbH",
            "Art. 6 Abs. 1 lit. c DSGVO",
            "10 Jahre",
            "Drittlaender",
        ],
        min_chars=120,
    ),
    "vvt_xlsx": ExtractionSample(
        key="vvt_xlsx",
        filename="vvt.xlsx",
        builder=_vvt_xlsx_bytes,
        expected_tokens=["Lohnabrechnung", "Art. 6 Abs. 1 lit. c DSGVO", "10 Jahre", "Bewerbermanagement"],
        column_header="| Zeile | A | B | C | D |",
        min_chars=60,
    ),
    "avv_docx": ExtractionSample(
        key="avv_docx",
        filename="avv.docx",
        builder=_avv_docx_bytes,
        expected_tokens=["Art. 28 DSGVO", "Gegenstand", "Lohnbuchhaltung", "Drittland", "Nein"],
        column_header="| Gegenstand | Lohnbuchhaltung |",
        min_chars=80,
    ),
    "tom_pptx": ExtractionSample(
        key="tom_pptx",
        filename="tom.pptx",
        builder=_tom_pptx_bytes,
        # Table content, slide-2 body, speaker notes AND the slide anchor must survive.
        expected_tokens=[
            "Verschlüsselung",
            "Zugriffskontrolle",
            "Löschkonzept",
            "AVV mit Hosting-Anbieter",
            "[Folie 2]",
        ],
        column_header="| Maßnahme | Status |",
        min_chars=60,
    ),
    "vvt_csv": ExtractionSample(
        key="vvt_csv",
        filename="vvt.csv",
        builder=_vvt_csv_bytes,
        expected_tokens=["Lohnabrechnung", "Art. 6 Abs. 1 lit. c DSGVO", "10 Jahre", "Bewerbermanagement"],
        column_header="| Zeile | A | B | C |",
        min_chars=60,
    ),
    "header_footer_docx": ExtractionSample(
        key="header_footer_docx",
        filename="hinweis.docx",
        builder=_header_footer_docx_bytes,
        # Controller (header) and version/file reference (footer) must survive extraction.
        expected_tokens=["Verantwortlicher: Muster GmbH", "Az. DS-2026-042", "Lohnabrechnung"],
        min_chars=40,
    ),
}
